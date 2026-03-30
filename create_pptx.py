from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG = RGBColor(0, 0, 0)
WHITE = RGBColor(0xF5, 0xF5, 0xF7)
MUTED = RGBColor(0x86, 0x86, 0x8B)
ACCENT = RGBColor(0x29, 0x97, 0xFF)
CARD_BG = RGBColor(0x1D, 0x1D, 0x1F)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

BLANK = prs.slide_layouts[6]  # blank layout


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER,
             font_name="Helvetica Neue", line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    return txBox, tf


def add_multirun_text(slide, runs, left, top, width, height,
                      alignment=PP_ALIGN.CENTER, line_spacing=1.15):
    """runs = list of (text, font_size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    for i, (text, size, color, bold) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Helvetica Neue"
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Subtle border
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


# ─────────────────────────────────────
# SLIDE 1 — The Problem
# ─────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
set_bg(s1)

add_text(s1, "THE PROBLEM", Inches(0), Inches(1.8), W, Inches(0.5),
         font_size=14, color=ACCENT, bold=True)

add_multirun_text(s1, [
    ("Brokers are always\n", 60, WHITE, True),
    ("one step behind.", 60, ACCENT, True),
], Inches(1.5), Inches(2.5), Inches(10.333), Inches(2))

add_text(s1, "Growth signals are scattered across dozens of sources.\nBy the time you spot an opportunity, a competitor has already made the call.",
         Inches(2.5), Inches(4.7), Inches(8.333), Inches(1.2),
         font_size=22, color=MUTED, bold=False, line_spacing=1.5)


# ─────────────────────────────────────
# SLIDE 2 — The Solution
# ─────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
set_bg(s2)

add_text(s2, "THE SOLUTION", Inches(0), Inches(1.8), W, Inches(0.5),
         font_size=14, color=ACCENT, bold=True)

add_multirun_text(s2, [
    ("An AI agent that finds\nyour next deal ", 60, WHITE, True),
    ("before\nanyone else.", 60, ACCENT, True),
], Inches(1.5), Inches(2.3), Inches(10.333), Inches(2.5))

add_text(s2, "It scrapes job boards, funding databases, and listings daily —\nthen delivers a scored prospect list every morning.",
         Inches(2.5), Inches(4.9), Inches(8.333), Inches(1.2),
         font_size=22, color=MUTED, bold=False, line_spacing=1.5)


# ─────────────────────────────────────
# SLIDE 3 — How It Works (Bento)
# ─────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
set_bg(s3)

add_text(s3, "HOW IT WORKS", Inches(0), Inches(0.7), W, Inches(0.5),
         font_size=14, color=ACCENT, bold=True)

add_multirun_text(s3, [
    ("Four steps. ", 44, WHITE, True),
    ("Fully automated.", 44, ACCENT, True),
], Inches(1.5), Inches(1.3), Inches(10.333), Inches(1))

cards = [
    ("📡", "Ingest", "Bright Data API pulls structured data\nfrom 13+ sources — job boards, company\ndatabases, property listings."),
    ("🧠", "Analyze", "AI extracts signals: hybrid job volume,\nheadcount growth, recent funding,\nsubmarket vacancy trends."),
    ("📊", "Score", "Weighted algorithm ranks every\nprospect 0–100 based on\nexpansion likelihood."),
    ("📬", "Deliver", "Dashboard with executive summary\ncards and a detailed spreadsheet —\ncontacts, trends, competitive density."),
]

card_w = Inches(4.8)
card_h = Inches(2.3)
gap = Inches(0.3)
grid_w = card_w * 2 + gap
start_x = (W - grid_w) // 2
start_y = Inches(2.6)

for i, (icon, title, desc) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = start_x + col * (card_w + gap)
    y = start_y + row * (card_h + gap)

    add_rounded_rect(s3, x, y, card_w, card_h)

    add_text(s3, icon, x + Inches(0.4), y + Inches(0.3), Inches(0.8), Inches(0.6),
             font_size=32, alignment=PP_ALIGN.LEFT)
    add_text(s3, title, x + Inches(0.4), y + Inches(0.9), card_w - Inches(0.8), Inches(0.4),
             font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    add_text(s3, desc, x + Inches(0.4), y + Inches(1.35), card_w - Inches(0.8), Inches(1.0),
             font_size=14, color=MUTED, bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.4)


# ─────────────────────────────────────
# SLIDE 4 — Impact
# ─────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
set_bg(s4)

add_text(s4, "IMPACT", Inches(0), Inches(1.4), W, Inches(0.5),
         font_size=14, color=ACCENT, bold=True)

add_multirun_text(s4, [
    ("Replace hours of research\nwith a ", 56, WHITE, True),
    ("5-minute briefing.", 56, ACCENT, True),
], Inches(1.5), Inches(2.0), Inches(10.333), Inches(2.2))

add_text(s4, "Multi-source intelligence catches prospects at the hiring and funding stage —\nmonths before they hit the broker market.",
         Inches(2.5), Inches(4.2), Inches(8.333), Inches(1.0),
         font_size=22, color=MUTED, bold=False, line_spacing=1.5)

# Stats row
stats = [
    ("13+", "Data Sources"),
    ("6–18 mo", "Earlier Signals"),
    ("5 min", "Morning Briefing"),
]

stat_w = Inches(3)
stats_total = stat_w * 3
stats_start = (W - stats_total) // 2
stat_y = Inches(5.5)

for i, (value, label) in enumerate(stats):
    x = stats_start + i * stat_w
    add_text(s4, value, x, stat_y, stat_w, Inches(0.8),
             font_size=48, color=ACCENT, bold=True)
    add_text(s4, label, x, stat_y + Inches(0.75), stat_w, Inches(0.4),
             font_size=14, color=MUTED, bold=False)


# Save
out = "/Users/lukadadiani/Documents/UM-hack/presentation.pptx"
prs.save(out)
print(f"Saved to {out}")
